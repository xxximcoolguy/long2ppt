"""longtext-to-slides · HTML → PPTX 导出器

策略：截图法 + 备注栏文本保留
- Playwright headless 启动 chromium，加载 HTML
- 每张 frame 切到 active 状态 → 截图 1920×1080 PNG
- python-pptx 构建 16:9 PPTX，每页一图填满
- 每页备注栏写入对应 frame 的主标题 + 主要文本（可搜可读）

依赖：
    pip install playwright python-pptx
    playwright install chromium

Usage:
    python export_pptx.py <input.html> [output.pptx]

默认输出：与 HTML 同目录，文件名替换扩展为 .pptx
"""
import sys
import os
import re
import time
import asyncio
import argparse
from pathlib import Path

try:
    from playwright.async_api import async_playwright
except ImportError:
    sys.exit("Missing dependency: playwright. Install:\n"
             "  pip install playwright python-pptx\n"
             "  playwright install chromium")

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
except ImportError:
    sys.exit("Missing dependency: python-pptx. Install: pip install python-pptx")


# ============================================================
# 1. 启动浏览器，捕获每张 frame 的截图 + 备注文本
# ============================================================

async def capture(html_path: Path, out_dir: Path,
                  vw: int = 1920, vh: int = 1080) -> list[dict]:
    """返回 [{index, image_path, title, notes}, ...]"""
    out_dir.mkdir(parents=True, exist_ok=True)
    url = f"file:///{str(html_path.resolve()).replace(os.sep, '/')}"

    captures: list[dict] = []
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        ctx = await browser.new_context(
            viewport={"width": vw, "height": vh},
            device_scale_factor=1.0,
        )
        page = await ctx.new_page()
        await page.goto(url, wait_until="networkidle")
        await page.wait_for_timeout(800)  # 等 deck.js 初始化

        # 拿到所有 frame 元素
        frame_count = await page.evaluate(
            "() => document.querySelectorAll('.deck .frame').length"
        )
        if frame_count == 0:
            await browser.close()
            raise SystemExit("[!] no .frame elements found in HTML")

        for i in range(frame_count):
            # 强制把第 i 张 frame 切到 active 状态
            await page.evaluate(f"""
                (() => {{
                    const fs = Array.from(document.querySelectorAll('.deck .frame'));
                    fs.forEach((f, j) => {{
                        if (j < {i}) f.dataset.state = 'past';
                        else if (j === {i}) f.dataset.state = 'active';
                        else f.dataset.state = 'next';
                    }});
                    /* 同步导航点 */
                    const pips = document.querySelectorAll('.rail .rail-pip');
                    pips.forEach((p, j) => p.setAttribute(
                        'aria-current', j === {i} ? 'true' : 'false'
                    ));
                    /* 隐藏 hotkey hint */
                    const h = document.querySelector('.hotkey-hint');
                    if (h) h.style.display = 'none';
                }})();
            """)
            # 等动画完成
            await page.wait_for_timeout(700)

            img_path = out_dir / f"frame-{i + 1:03d}.png"
            await page.screenshot(path=str(img_path), full_page=False, type="png")

            # 提取 frame 文本 (主标题 + 内容文字) 作为备注
            data = await page.evaluate(f"""
                (() => {{
                    const f = document.querySelectorAll('.deck .frame')[{i}];
                    if (!f) return {{ title: '', notes: '' }};
                    const title = f.dataset.title || '';
                    /* 抽取主要标题与正文 */
                    const lines = [];
                    f.querySelectorAll('h1, h2, h3, h4, .lede, .cue, .deck-meta, p, li').forEach(n => {{
                        const t = (n.innerText || '').trim();
                        if (t && t.length > 0 && t.length < 600) lines.push(t);
                    }});
                    const seen = new Set();
                    const unique = [];
                    lines.forEach(l => {{ if (!seen.has(l)) {{ seen.add(l); unique.push(l); }} }});
                    return {{ title, notes: unique.join('\\n') }};
                }})();
            """)

            captures.append({
                "index": i + 1,
                "image_path": img_path,
                "title": data.get("title", "") or f"Frame {i + 1}",
                "notes": data.get("notes", ""),
            })

            print(f"[*] captured frame {i+1}/{frame_count}: {data.get('title','')}", file=sys.stderr)

        await browser.close()
    return captures


# ============================================================
# 2. 用 python-pptx 拼装 PPTX
# ============================================================

def build_pptx(captures: list[dict], out_path: Path,
               slide_w_inches: float = 13.333,
               slide_h_inches: float = 7.5) -> None:
    """slide_w/h 默认 16:9（13.333 × 7.5 英寸）"""
    prs = Presentation()
    prs.slide_width = Inches(slide_w_inches)
    prs.slide_height = Inches(slide_h_inches)

    blank_layout = prs.slide_layouts[6]  # 空白布局

    for c in captures:
        slide = prs.slides.add_slide(blank_layout)
        # 图片填满
        slide.shapes.add_picture(
            str(c["image_path"]),
            left=Emu(0), top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
        # 备注栏写大纲
        notes = slide.notes_slide.notes_text_frame
        notes.text = f"【{c['title']}】\n\n{c['notes']}"

    prs.save(str(out_path))


# ============================================================
# Entry
# ============================================================

async def main(args):
    html_path = Path(args.html).resolve()
    if not html_path.exists():
        sys.exit(f"[!] HTML not found: {html_path}")

    out_pptx = Path(args.output) if args.output else html_path.with_suffix(".pptx")
    tmp_dir = out_pptx.parent / f"{out_pptx.stem}-frames"

    print(f"[*] capturing frames @ {args.width}x{args.height} → {tmp_dir}/", file=sys.stderr)
    captures = await capture(html_path, tmp_dir, vw=args.width, vh=args.height)

    print(f"[*] building pptx → {out_pptx}", file=sys.stderr)
    build_pptx(captures, out_pptx)

    if not args.keep_frames:
        for c in captures:
            try: c["image_path"].unlink()
            except: pass
        try: tmp_dir.rmdir()
        except: pass

    print(f"OK  frames={len(captures)}  out={out_pptx}")


def parse_args():
    ap = argparse.ArgumentParser(description="Convert longtext-to-slides HTML to PPTX")
    ap.add_argument("html", help="Path to the HTML deck")
    ap.add_argument("output", nargs="?", default=None, help="Output .pptx (default: same name as HTML)")
    ap.add_argument("--width", type=int, default=1920, help="Capture viewport width (default 1920)")
    ap.add_argument("--height", type=int, default=1080, help="Capture viewport height (default 1080)")
    ap.add_argument("--keep-frames", action="store_true", help="Keep extracted PNG frames after build")
    return ap.parse_args()


if __name__ == "__main__":
    args = parse_args()
    asyncio.run(main(args))
